from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
PURPLE = RGBColor(124, 58, 237)
GOLD = RGBColor(245, 158, 11)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
DARK_GRAY = RGBColor(100, 100, 100)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)

def add_title_slide(title, subtitle, period):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    period_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.8))
    tf = period_box.text_frame
    p = tf.paragraphs[0]
    p.text = period
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, items, highlights=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    
    y_pos = 1.8
    for i, item in enumerate(items):
        item_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.7))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        
        if highlights and i in highlights:
            p.font.color.rgb = GOLD
            p.font.bold = True
        
        y_pos += 0.75
    
    return slide

def add_table_slide(title, data, columns):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    rows = len(data) + 1
    cols = len(columns)
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.333)
    height = Inches(0.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if row_idx % 2 == 0 else RGBColor(30, 41, 59)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# Slide 1: Title
add_title_slide(
    "Analyse Commerciale Avril 2026",
    "Chiffre dAffaires VC - SMG",
    "Periode: 01-21 Avril 2026 (N) vs 2025 (N-1)"
)

# Slide 2: Analyse Globale
add_content_slide(
    "Analyse Globale",
    [
        "CA Total N: 914 037 EUR (+25.4%)",
        "CA Total N-1: 729 166 EUR",
        "Nombre de factures: 533 (vs 477, +11.7%)",
        "Panier moyen: 1 715 EUR (vs 1 529 EUR, +12.2%)",
        "",
        "Tendance: CROISSANCE FORTE",
        "Facteurs: Effet convention, saisonnalite, hausse trafic"
    ],
    [0, 1, 5]
)

# Slide 3: Top 5 Magasins
add_table_slide(
    "Top 5 - Meilleures Performances",
    [
        ["111", "176 788 EUR", "+109 916 EUR", "+164.4%"],
        ["116", "77 354 EUR", "+22 108 EUR", "+40.0%"],
        ["311", "44 492 EUR", "+1 949 EUR", "+4.6%"],
        ["409", "26 175 EUR", "+1 800 EUR", "+7.4%"],
        ["140", "21 845 EUR", "+2 000 EUR", "+10.1%"]
    ],
    ["Magasin", "CA N", "Ecart", "Evolution"]
)

# Slide 4: Flop 5 Magasins
add_table_slide(
    "Flop 5 - Plus Fortes Baisses",
    [
        ["143", "8 331 EUR", "-57 429 EUR", "-87.3%"],
        ["511", "39 406 EUR", "-23 904 EUR", "-37.8%"],
        ["313", "27 309 EUR", "-16 262 EUR", "-37.3%"],
        ["524", "10 183 EUR", "-7 966 EUR", "-43.9%"],
        ["215", "47 005 EUR", "-14 014 EUR", "-23.0%"]
    ],
    ["Magasin", "CA N", "Ecart", "Evolution"]
)

# Slide 5: Analyse par Convention
add_table_slide(
    "Analyse par Convention",
    [
        ["VC.CONV.", "550 407 EUR", "528 994 EUR", "+21 413 EUR", "+4.0%"],
        ["VC.CONSO.", "345 529 EUR", "197 722 EUR", "+147 806 EUR", "+74.8%"],
        ["VC.PARTIC.", "18 102 EUR", "2 449 EUR", "+15 653 EUR", "+639.1%"]
    ],
    ["Convention", "CA N", "CA N-1", "Ecart", "Evolution"]
)

# Slide 6: Synthese
add_content_slide(
    "Synthese - 5 Enseignements Cles",
    [
        "1. Croissance globale exceptionnelle: +25.4% (185k EUR)",
        "2. Magasin 111 star: +110k EUR (+165%) - 19% du CA total",
        "3. Polarisation reseau: ecart fort entre top et flop",
        "4. VC.CONSO moteur: +75% porte la croissance",
        "5. Panier moyen en hausse: +12% - monte en gamme"
    ],
    [0, 1, 3]
)

# Slide 7: Risques
add_content_slide(
    "Risques Identifies",
    [
        "Dependance Magasin 111: -19% du CA si perte convention",
        "Surexposition VC.CONSO: Risque credit en cas de retournement",
        "Effondrement Magasin 143: -87% soit -57k EUR",
        "Stagnation VC.CONV: +4% seulement vs potentiel +10%",
        "Baisses en cascade (511, 313): -68k EUR cumules"
    ],
    [0, 1, 2]
)

# Slide 8: Plan Action
add_content_slide(
    "Plan dAction Recommande",
    [
        "URGENT (0-2 semaines):",
        "  - Diagnostic urgence Magasin 143 (recuperation ~30k EUR)",
        "  - Consolider convention Magasin 111 (+20k EUR)",
        "  - Relancer partenaires VC.CONV. inactifs (+15k EUR)",
        "",
        "PRIORITAIRE (avril-mai):",
        "  - Plan action Magasins 511/313",
        "  - Suivi hebdomadaire VC.CONSO",
        "",
        "STRATEGIQUE (Q2-Q3):",
        "  - Reduire dependence VC.CONSO",
        "  - Programme montee en gamme reseau"
    ],
    [0, 3, 6]
)

# Save
output_path = r"C:\Users\hachk\pilotage_b2b\__pycache__\Rapport_CA_Avril_2026.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
