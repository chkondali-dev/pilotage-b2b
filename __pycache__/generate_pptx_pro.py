from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsmap

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
PURPLE = RGBColor(124, 58, 237)
GOLD = RGBColor(245, 158, 11)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(148, 163, 184)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)

def create_gradient_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    return bg

def add_title(text, slide, size=44):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return title_box

def add_subtitle(text, slide):
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(2), Inches(0.08))
    sub_box.fill.solid()
    sub_box.fill.fore_color.rgb = GOLD
    sub_box.line.fill.background()

# Slide 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide1)
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Analyse Commerciale Avril 2026"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
sub_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Chiffre dAffaires VC - SMG"
p.font.size = Pt(32)
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER
period_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.8))
tf = period_box.text_frame
p = tf.paragraphs[0]
p.text = "Periode: 01-21 Avril 2026 (N) vs 2025 (N-1)"
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 2: Analysis Globale
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide2)
add_title("Analyse Globale", slide2)
add_subtitle(None, slide2)

metrics = [
    ("914 037 EUR", "+25.4%", "CA Total"),
    ("533", "+11.7%", "Nb Factures"),
    ("1 715 EUR", "+12.2%", "Panier Moyen")
]
x_pos = 0.8
for val, pct, label in metrics:
    box = slide2.shapes.add_textbox(Inches(x_pos), Inches(1.8), Inches(3.8), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = pct
    p2.font.size = Pt(24)
    p2.font.color.rgb = GREEN
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = label
    p3.font.size = Pt(16)
    p3.font.color.rgb = LIGHT_GRAY
    p3.alignment = PP_ALIGN.CENTER
    x_pos += 4.1

# Slide 3: Top 5
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide3)
add_title("Top 5 - Meilleures Performances", slide3)

top_data = [
    ("MAGASIN", "CA N", "EVOLUTION"),
    ("111", "176 788 EUR", "+164.4%"),
    ("116", "77 354 EUR", "+40.0%"),
    ("311", "44 492 EUR", "+4.6%"),
    ("409", "26 175 EUR", "+7.4%"),
    ("140", "21 845 EUR", "+10.1%")
]
rows = len(top_data)
cols = 3
table3 = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5)).table
for i, row in enumerate(top_data):
    for j, val in enumerate(row):
        cell = table3.cell(i, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE if i == 0 else NAVY
        cell.text_frame.paragraphs[0].font.size = Pt(16) if i == 0 else Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True if i == 0 else False
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 4: Flop 5
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide4)
add_title("Flop 5 - Plus Fortes Baisses", slide4)

flop_data = [
    ("MAGASIN", "CA N", "EVOLUTION"),
    ("143", "8 331 EUR", "-87.3%"),
    ("511", "39 406 EUR", "-37.8%"),
    ("313", "27 309 EUR", "-37.3%"),
    ("524", "10 183 EUR", "-43.9%"),
    ("215", "47 005 EUR", "-23.0%")
]
rows = len(flop_data)
table4 = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5)).table
for i, row in enumerate(flop_data):
    for j, val in enumerate(row):
        cell = table4.cell(i, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED if i == 0 else NAVY
        cell.text_frame.paragraphs[0].font.size = Pt(16) if i == 0 else Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True if i == 0 else False
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 5: Conventions
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide5)
add_title("Analyse par Convention", slide5)

conv_data = [
    ("CONVENTION", "CA N", "CA N-1", "EVOLUTION"),
    ("VC.CONV.", "550 407 EUR", "528 994 EUR", "+4.0%"),
    ("VC.CONSO.", "345 529 EUR", "197 722 EUR", "+74.8%"),
    ("VC.PARTIC.", "18 102 EUR", "2 449 EUR", "+639%")
]
rows = len(conv_data)
cols = 4
table5 = slide5.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5)).table
for i, row in enumerate(conv_data):
    for j, val in enumerate(row):
        cell = table5.cell(i, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE if i == 0 else NAVY
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True if i == 0 else False
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 6: Synthese
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide6)
add_title("Synthese - 5 Enseignements Cles", slide6)

lessons = [
    "1. Croissance globale +25.4% - Performance exceptionnelle",
    "2. Magasin 111 star: +110kEUR (+165%) - 19% du CA total",
    "3. Polarisation reseau: ecart fort entre top et flop",
    "4. VC.CONSO moteur: +75% porte la croissance",
    "5. Panier moyen en hausse: +12% - montee en gamme"
]
y = 1.6
for lesson in lessons:
    box = slide6.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = lesson
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    y += 0.85

# Slide 7: Risques
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide7)
add_title("Risques Identifies", slide7)

risques = [
    "Dependance Magasin 111: -19% du CA si perte convention",
    "Surexposition VC.CONSO: Risque credit en cas de retournement",
    "Effondrement Magasin 143: -87% soit -57kEUR",
    "Stagnation VC.CONV: +4% seulement vs potentiel +10%",
    "Baisses en cascade (511, 313): -68kEUR cumules"
]
y = 1.6
for risque in risques:
    box = slide7.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = risque
    p.font.size = Pt(20)
    p.font.color.rgb = RED if "143" in risque else WHITE
    y += 0.75

# Slide 8: Plan Action
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
create_gradient_background(slide8)
add_title("Plan dAction Recommande", slide8)

actions = [
    "URGENT (0-2 semaines):",
    "  - Diagnostic urgence Magasin 143 (recuperation ~30kEUR)",
    "  - Consolider convention Magasin 111 (+20kEUR)",
    "  - Relancer partenaires VC.CONV. inactifs (+15kEUR)",
    "",
    "PRIORITAIRE (avril-mai):",
    "  - Plan action Magasins 511/313",
    "  - Suivi hebdomadaire VC.CONSO",
    "",
    "STRATEGIQUE (Q2-Q3):",
    "  - Reduire dependence VC.CONSO",
    "  - Programme montee en gamme reseau"
]
y = 1.4
for action in actions:
    box = slide8.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = action
    p.font.size = Pt(18)
    p.font.bold = True if action.endswith(":") else False
    p.font.color.rgb = GOLD if action.endswith(":") else WHITE
    y += 0.55

output_path = r"C:\Users\hachk\pilotage_b2b\__pycache__\Rapport_CA_Avril_2026_v2.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")